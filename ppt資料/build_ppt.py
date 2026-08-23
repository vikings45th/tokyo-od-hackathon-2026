# -*- coding: utf-8 -*-
"""都知事杯オープンデータ・ハッカソン2026 — 2分プレゼン資料（5枚）を生成する。

  1. 共働き世帯の課題の中での位置づけ（漏斗）＋ 規模の実測
  2. 構造的課題＝「小1の壁」の正体（増える需要／増えない受け皿／出ない数字）
  3. サービスとデモ
  4. データとロジック（何を結合し、どこまで測り、どこを測っていないか）
  5. 展望（階段）— 解ける課題が広がる。それを支えるのは運用コストがゼロなこと

1枚目の漏斗と5枚目の階段は同じ木を上下から見たもの。
「機能を足す」ではなく「最初に見せた木を、この順で登る」にするための構造。

参考にしたレイアウト：ppt資料/ の戦略コンサル2本
  アクションタイトル／2パネル＋細い縦罫／コールアウト＋引き出し線／出典行。
  密度は真似しない。あれは委員が手元で読む資料、こちらは2分動画で「見る」もの。

🔴 数字はすべて data/app/data.json（都の公開データ）と src/core の実測値。
   測っていないことを測ったように書かないこと。

フォントは Meiryo UI（欧文 latin ／ 和文 ea ／ 記号 cs のすべて）。
サービス名を変えるときは SERVICE を書き換えて再生成する。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.oxml.ns import qn

SERVICE = '小1のまち'
FONT = 'Meiryo UI'

# src/ui/palette.ts の実配色に合わせる（スライドと実画面の色を一致させる）
ACCENT      = RGBColor(0x1C, 0x5C, 0xAB)
ACCENT_DEEP = RGBColor(0x0D, 0x36, 0x6B)
ACCENT_MID  = RGBColor(0x55, 0x98, 0xE7)
ACCENT_PALE = RGBColor(0x86, 0xB6, 0xEF)
TINT        = RGBColor(0xEC, 0xF2, 0xFB)
TINT2       = RGBColor(0xF4, 0xF7, 0xFC)
INK         = RGBColor(0x1A, 0x1A, 0x1A)
GREY        = RGBColor(0x5A, 0x5A, 0x5A)
GREY_LT     = RGBColor(0x9A, 0x9A, 0x9A)
RULE        = RGBColor(0xD5, 0xD5, 0xD1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = 13.333, 7.5
L, R = 0.62, 0.62
CW = SW - L - R

SOURCE = ('出典：東京都福祉局「東京の学童クラブ事業実施状況」／東京都教育庁「教育人口等推計」／'
          '東京都教育庁「公立学校統計調査報告書（東京都公立学校一覧）」'
          '（いずれも東京都オープンデータカタログ・CC BY 4.0）')

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
BLANK = prs.slide_layouts[6]


# ── 部品 ──────────────────────────────────────────────
def tb(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return box, tf


def put(tf, text, size, color=INK, bold=False, space_after=0, line_sp=None, align=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.space_after = Pt(space_after)
    if line_sp:
        p.line_spacing = line_sp
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return p


def lines(tf, texts, size, color=INK, line_sp=1.35, space_after=0, bold=False, align=None):
    for i, t in enumerate(texts):
        put(tf, t, size, color, bold=bold, line_sp=line_sp,
            space_after=space_after, align=align, first=(i == 0))


def rect(slide, x, y, w, h, fill=None, outline=None, lw=1.0, dash=False, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if outline is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = outline
        s.line.width = Pt(lw)
        if dash:
            ln = s.line._get_or_add_ln()
            ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    tf = s.text_frame
    tf.word_wrap = True
    # 小さい図形（丸数字・チップ・0.055インチの罫）で余白が効くと文字が折り返す
    tf.margin_left = tf.margin_right = Inches(0.14 if w > 1.2 else 0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.06 if h > 0.5 else 0.0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return s


def line(slide, x1, y1, x2, y2, color=RULE, lw=1.0, arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(lw)
    if arrow:
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def title(slide, text, sub=None, size=26):
    _, tf = tb(slide, L, 0.40, CW, 1.05)
    put(tf, text, size, ACCENT, bold=True, line_sp=1.22, first=True)
    if sub:
        _, tf2 = tb(slide, L, 1.34, CW, 0.40)
        put(tf2, sub, 12.5, GREY, line_sp=1.3, first=True)


def footer(slide, n, source=SOURCE):
    _, tf = tb(slide, L, 6.98, CW - 0.9, 0.42)
    put(tf, source, 7.5, GREY_LT, line_sp=1.25, first=True)
    _, tf2 = tb(slide, SW - R - 0.6, 6.98, 0.6, 0.25, align=PP_ALIGN.RIGHT)
    put(tf2, str(n), 9, GREY_LT, align=PP_ALIGN.RIGHT, first=True)


def placeholder(slide, x, y, w, h, label, note):
    """実画面のキャプチャを貼る枠。貼ったらこの図形は削除する。"""
    s = rect(slide, x, y, w, h, fill=TINT, outline=ACCENT_MID, lw=1.25, dash=True)
    put(s.text_frame, label, 15, ACCENT, bold=True, align=PP_ALIGN.CENTER, space_after=5, first=True)
    put(s.text_frame, note, 10.5, GREY, align=PP_ALIGN.CENTER, line_sp=1.4)
    return s


def band(slide, x, y, w, text, fill=ACCENT):
    s = rect(slide, x, y, w, 0.36, fill=fill)
    put(s.text_frame, text, 12, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    return s


def kicker(slide, y, text, size=16):
    s = rect(slide, L, y, CW, 0.62, fill=TINT)
    put(s.text_frame, text, size, ACCENT_DEEP, bold=True, align=PP_ALIGN.CENTER, first=True)
    return s


def timeline(slide, YL, num_size=18, desc_size=13):
    """家を買う年と、学童に入れない年の6年のズレ。"""
    line(slide, 1.55, YL, 11.75, YL, color=RULE, lw=2.0)
    for cx, col, year, desc, dark in [
        (2.45, GREY_LT, '2026年', '第一子が0歳。住宅ローンを組む', False),
        (10.85, ACCENT, '2032年 4月', '子が小1。学童に入れない', True),
    ]:
        d = 0.28
        rect(slide, cx - d / 2, YL - d / 2, d, d, fill=col, shape=MSO_SHAPE.OVAL)
        _, tf = tb(slide, cx - 1.9, YL - 0.82, 3.8, 0.40,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        put(tf, year, num_size, ACCENT_DEEP if dark else GREY, bold=True,
            align=PP_ALIGN.CENTER, first=True)
        _, tf = tb(slide, cx - 1.9, YL + 0.26, 3.8, 0.4, align=PP_ALIGN.CENTER)
        put(tf, desc, desc_size, INK if dark else GREY, bold=dark,
            align=PP_ALIGN.CENTER, first=True)
    line(slide, 4.35, YL, 10.58, YL, color=ACCENT, lw=2.5, arrow=True)
    g = rect(slide, 6.10, YL - 0.25, 1.15, 0.50, fill=WHITE)
    put(g.text_frame, '6年', 21, ACCENT, bold=True, align=PP_ALIGN.CENTER, first=True)


# ══ 1枚目：課題の中での位置づけ（漏斗）＋ 規模 ═════════
s1 = prs.slides.add_slide(BLANK)
title(s1, '共働き世帯の両立を阻むものが、保育園から小学校に移った',
      '東京49自治体・2025年5月1日時点の実測と、受け皿を据え置いた場合の試算')

# 左：3段の漏斗
_, tf = tb(s1, L, 1.92, 5.5, 0.28)
put(tf, '共働き・子育て世帯が抱える課題', 12, GREY, bold=True, first=True)

funnel = [
    (L,        5.50, TINT2, GREY,       '仕事と育児の両立', 15, False),
    (L + 0.40, 4.70, TINT,  ACCENT_DEEP, '小学校に上がる時点の断絶（「小1の壁」）', 14, False),
    (L + 0.80, 3.90, ACCENT, WHITE,      'その街で、学童に入れるか　← v1 が解く', 14, True),
]
for i, (x, w, fill, col, txt, sz, bold) in enumerate(funnel):
    y = 2.28 + i * 0.78
    s = rect(s1, x, y, w, 0.60, fill=fill)
    put(s.text_frame, txt, sz, col, bold=True, first=True)
    if i < 2:
        line(s1, x + 0.55, y + 0.60, x + 0.55, y + 0.78, color=ACCENT_MID, lw=1.75, arrow=True)

_, tf = tb(s1, L, 4.78, 5.5, 0.62)
lines(tf, ['保育所の待機は都の施策で大きく前進した。',
           '次に来るのが、小学校に上がる時点の断絶'], 11.5, GREY, line_sp=1.35)

line(s1, 6.42, 1.92, 6.42, 5.42, color=RULE, lw=1.0)

# 右：規模の実測3数字
_, tf = tb(s1, 6.72, 1.92, 5.99, 0.28)
put(tf, 'その規模（都の公開データで実測）', 12, GREY, bold=True, first=True)

tiles = [
    ('146,393人', '学童に登録している児童。公立小学校児童 589,912人の 4人に1人',
     '2025年5月1日・49自治体'),
    ('+5.9%', '2年で増えた学童の需要（登録＋待機）。同じ2年で児童数は −0.8%',
     '2023年5月→2025年5月・48自治体'),
    ('19,573人分', '受け皿が今のままなら、2038年度に足りない',
     '本サービスの試算・46/49自治体で不足'),
]
for i, (num, desc, note) in enumerate(tiles):
    y = 2.28 + i * 1.08
    rect(s1, 6.72, y, 0.055, 0.92, fill=ACCENT)
    rect(s1, 6.775, y, 5.935, 0.92, fill=TINT)
    _, tf = tb(s1, 6.95, y + 0.10, 2.15, 0.42)
    put(tf, num, 24, ACCENT, bold=True, first=True)
    _, tf = tb(s1, 9.20, y + 0.10, 3.35, 0.72)
    put(tf, desc, 10.5, INK, line_sp=1.3, space_after=2, first=True)
    put(tf, note, 8.5, GREY_LT, line_sp=1.2)

kicker(s1, 5.72, '保育園の壁は下がった。壁が消えたのではなく、小学校に移っただけ。')
footer(s1, 1, '出典：東京都福祉局「東京の学童クラブ事業実施状況」／東京都教育庁「教育人口等推計」'
              '（東京都オープンデータカタログ・CC BY 4.0）。'
              '需要の伸びは、計上方法が変わった江戸川区を除く48自治体で実測')


# ══ 2枚目：構造的課題＝「小1の壁」の正体 ═══════════════
s2 = prs.slides.add_slide(BLANK)
title(s2, '「小1の壁」の正体は、3つが同時に起きていること',
      '需要は増える／受け皿は増えない／それが数字に出ない。'
      'だから「今年の待機児童数」を見ても、自分の街の6年後は分からない')

CWD = (CW - 0.60) / 3
cols = [
    ('1', '需要は増えている',
     ['共働きが増え、学童を使う子の割合は',
      '年 +0.81pt で上がり続けている。',
      '登録＋待機は2年で +5.9%'],
     '48自治体・2023-05→2025-05 の実測'),
    ('2', '受け皿は増えていない',
     ['2年でクラブが1つも増えていない自治体が',
      '25／49。うち6自治体は減っている',
      '（都計の増加は江戸川区がほぼ全部）'],
     '49自治体・クラブ数の実測'),
    ('3', 'それが数字に出ない',
     ['待機児童は49自治体のうち19がゼロ。',
      '申込前に断念した人は、登録にも',
      '待機にも計上されない'],
     '2025-05-01 の実測'),
]
for i, (num, head, body, note) in enumerate(cols):
    x = L + i * (CWD + 0.30)
    rect(s2, x, 2.15, CWD, 1.72, fill=TINT)
    n = rect(s2, x + 0.22, 2.34, 0.34, 0.34, fill=ACCENT, shape=MSO_SHAPE.OVAL)
    put(n.text_frame, num, 12.5, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    _, tf = tb(s2, x + 0.66, 2.38, CWD - 0.88, 0.34)
    put(tf, head, 15, ACCENT_DEEP, bold=True, first=True)
    _, tf = tb(s2, x + 0.22, 2.86, CWD - 0.44, 0.72)
    lines(tf, body, 11, INK, line_sp=1.32)
    _, tf = tb(s2, x + 0.22, 3.54, CWD - 0.44, 0.24)
    put(tf, note, 9, GREY_LT, first=True)

_, tf = tb(s2, L, 4.08, CW, 0.34, align=PP_ALIGN.CENTER)
put(tf, 'だから、家を買う時点では、その街の6年後が分からない',
    16, ACCENT_DEEP, bold=True, align=PP_ALIGN.CENTER, first=True)

timeline(s2, 5.30, num_size=17, desc_size=12.5)
footer(s2, 2, SOURCE + '。需要の伸びは src/core/trend.ts、受け皿と待機は data/app/data.json で再現できる')


# ══ 3枚目：サービスとデモ ════════════════════════════
s3 = prs.slides.add_slide(BLANK)
title(s3, '生まれ年を入れるだけで、その子が小1になる年の東京が出る')

_, tf = tb(s3, L, 1.30, 4.0, 0.55)
put(tf, SERVICE, 30, ACCENT, bold=True, first=True)
_, tf = tb(s3, L + 2.35, 1.44, 8.5, 0.4)
put(tf, '東京49自治体 × 入学年度2027〜2038 ＝ 588通りを、1画面で比べる', 13, GREY, first=True)

placeholder(s3, L, 2.00, CW, 4.40,
            '▶ ここに実画面のキャプチャ／画面録画（全画面）',
            '生まれ年を入れる → 地図が入学年度に切り替わる → 中央区を選ぶ →'
            '「882人分、足りない」→ 隣接区の比較。\n'
            '動画では40秒。枠は 12.09 × 4.40 inch（比率 2.75 : 1）')

co = rect(s3, 8.45, 2.55, 3.55, 1.30, fill=WHITE, outline=ACCENT, lw=1.75)
put(co.text_frame, '2031年度・中央区', 11, GREY, align=PP_ALIGN.CENTER, space_after=3, first=True)
put(co.text_frame, '882人分、足りない', 20, ACCENT, bold=True, align=PP_ALIGN.CENTER, space_after=3)
put(co.text_frame, '多く見れば 1,543人', 11, GREY, align=PP_ALIGN.CENTER)
line(s3, 8.45, 3.20, 7.20, 3.90, color=ACCENT, lw=1.5)

kicker(s3, 6.55, '「今年の待機児童数」ではなく、「あなたの子が小1になる年」で並べ替える。', size=15)
footer(s3, 3)


# ══ 4枚目：データとロジック ═══════════════════════════
s4 = prs.slides.add_slide(BLANK)
title(s4, '都の公式データ3種を結合し、予測が外れる幅まで測っている',
      '49自治体 × 12年度 ＝ 588セル。サーバーもDBもAPIキーも使わず、すべてブラウザの中で計算する')

FW = (CW - 3 * 0.42) / 4
flow = [
    ('3つの公式データ', ['学童クラブ実施状況', '教育人口等推計（3世代）', '公立学校一覧 1,241校']),
    ('結合して整える', ['ヘッダ6行・区市別2ブロックの', '年次CSVを正規化し、',
                        '49自治体×12年度に揃える']),
    ('予測する', ['都の公式推計を読み、切れる', '2031年度以降は全都の伸び率で接続。',
                  '登録率は最小二乗＋縮約推定']),
    ('外れ幅を測る', ['推計3世代を突き合わせ、', '実際に何%外れたかを算出。',
                      '帯はその実測から引く']),
]
for i, (head, body) in enumerate(flow):
    x = L + i * (FW + 0.42)
    rect(s4, x, 1.98, FW, 0.055, fill=ACCENT)
    rect(s4, x, 2.035, FW, 1.28, fill=TINT)
    _, tf = tb(s4, x + 0.18, 2.16, FW - 0.36, 0.32)
    put(tf, head, 13.5, ACCENT_DEEP, bold=True, first=True)
    _, tf = tb(s4, x + 0.18, 2.54, FW - 0.36, 0.70)
    lines(tf, body, 9.5, INK, line_sp=1.32)
    if i < 3:
        line(s4, x + FW + 0.09, 2.68, x + FW + 0.34, 2.68, color=ACCENT, lw=2.0, arrow=True)

PT = 3.60
band(s4, L, PT, 5.55, '予測が外れた幅（実測・絶対誤差平均）')
cd = CategoryChartData()
cd.categories = ['1年先', '2年先']
cd.add_series('誤差', (0.93, 1.37))
gf = s4.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                         Inches(L), Inches(PT + 0.42), Inches(5.55), Inches(2.05), cd)
ch = gf.chart
ch.has_legend = False
ch.has_title = False
ch.font.size = Pt(12)
ch.font.name = FONT
ch.font.color.rgb = INK
pl = ch.plots[0]
pl.gap_width = 150
pl.has_data_labels = True
dl = pl.data_labels
dl.number_format = '0.00"%"'
dl.number_format_is_linked = False
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.font.size = Pt(17)
dl.font.bold = True
dl.font.name = FONT
dl.font.color.rgb = ACCENT_DEEP
ch.value_axis.has_major_gridlines = False
ch.value_axis.visible = False
ch.category_axis.has_major_gridlines = False
ch.category_axis.format.line.color.rgb = RULE
pts = pl.series[0].points
pts[0].format.fill.solid(); pts[0].format.fill.fore_color.rgb = ACCENT_PALE
pts[1].format.fill.solid(); pts[1].format.fill.fore_color.rgb = ACCENT

_, tf = tb(s4, L, PT + 2.52, 5.55, 0.4)
put(tf, '令和5・6年度版が出した推計を、令和7年度の実数と突き合わせた結果（49自治体）',
    10, GREY, line_sp=1.35, first=True)

line(s4, 6.44, PT, 6.44, PT + 2.92, color=RULE, lw=1.0)

band(s4, 6.72, PT, 5.99, '測っていないことは、測っていないと書く')
_, tf = tb(s4, 6.72, PT + 0.54, 5.99, 2.20)
for i, (h, b) in enumerate([
    ('2031年度から先', '都の公式推計が切れるので「接続した推定」と画面に明記し、帯を1.5倍に広げる'),
    ('自治体別トレンド', '推定する仕組みは実装済み。ただし時点が2つしかない今は'
                        '都全体の実測値にフォールバックしている、と画面に出す'),
    ('データがない自治体', '0点ではなく灰色で描く。「足りている」と「分からない」を同じ色にしない'),
]):
    put(tf, h, 11, ACCENT_DEEP, bold=True, space_after=1, first=(i == 0))
    put(tf, b, 10, INK, line_sp=1.32, space_after=7)
footer(s4, 4, SOURCE + '。バックテストは src/core/backtest.ts と npm test で再現できる')


# ══ 5枚目：展望（階段）＋ それを支える土台 ═════════════
s5 = prs.slides.add_slide(BLANK)
title(s5, '同じ器のまま、解ける課題を上に広げていく',
      '1枚目の木を、この順で登る。器（Indicator 配列）は変えず、軸を足すだけ')

SBW = 3.81
steps = [
    ('STEP 1', 'いま（v1）', 'その街で、学童に入れるか',
     ['学童の受け皿リスクを', '49自治体 × 12年度で比べる']),
    ('STEP 2', '次', 'その街で、両立を続けられるか',
     ['通勤時間・保育園の空き・学区を', '同じ器に軸として足す']),
    ('STEP 3', 'その先', 'どの街に、どれだけ受け皿が要るか',
     ['使い手が住民から自治体へ。', '整備計画の需給見通しに同じ数字を']),
]
for i, (chip, when, q, body) in enumerate(steps):
    x = L + i * (SBW + 0.33)
    y = 3.66 - i * 0.76
    accent_col = ACCENT if i == 0 else ACCENT_MID
    rect(s5, x, y, SBW, 0.055, fill=accent_col)
    rect(s5, x, y + 0.055, SBW, 1.42, fill=TINT if i == 0 else TINT2)
    c = rect(s5, x + 0.20, y + 0.18, 0.86, 0.26, fill=accent_col)
    put(c.text_frame, chip, 9, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    _, tf = tb(s5, x + 1.14, y + 0.20, SBW - 1.34, 0.24)
    put(tf, when, 10, GREY, first=True)
    _, tf = tb(s5, x + 0.20, y + 0.52, SBW - 0.40, 0.36)
    put(tf, q, 14, ACCENT_DEEP, bold=True, line_sp=1.2, first=True)
    _, tf = tb(s5, x + 0.20, y + 0.96, SBW - 0.40, 0.44)
    lines(tf, body, 10, INK, line_sp=1.3)
    if i < 2:
        line(s5, x + SBW + 0.04, y + 0.40, x + SBW + 0.29, y - 0.30,
             color=ACCENT_MID, lw=2.0, arrow=True)

# 階段を支える土台
rect(s5, L, 5.32, CW, 0.055, fill=ACCENT_DEEP)
base = rect(s5, L, 5.375, CW, 1.10, fill=TINT)
_, tf = tb(s5, L + 0.28, 5.50, 4.4, 0.34)
put(tf, 'この階段を登れる理由', 14, ACCENT_DEEP, bold=True, first=True)
_, tf = tb(s5, L + 0.28, 5.88, 5.6, 0.50)
lines(tf, ['使うのは毎年更新される都の公式データだけ。',
           'APIキーもDBもサーバーも持たない'], 10.5, INK, line_sp=1.32)
line(s5, 6.60, 5.50, 6.60, 6.30, color=RULE, lw=1.0)
_, tf = tb(s5, 6.85, 5.50, 5.60, 0.90)
lines(tf, ['年1回データを差し替えれば動き続ける（静的配信 501KB／gzip 161KB）。',
           '軸を足しても画面もヒートマップも変更が要らない（要件 FR-9）。',
           '運用コストは実質ゼロ。作った人が抜けても止まらない。'],
      10.5, INK, line_sp=1.32)

kicker(s5, 6.60, '一年生になったら、その街はどうなっているか。先に、見にいける。', size=15)
footer(s5, 5, '出典：東京都オープンデータカタログ（CC BY 4.0）。'
              '配信サイズは npm run build の実測値')


# ── 全ランに Meiryo UI（latin / ea / cs）を焼き込む ──
def stamp(el):
    for rPr in el.iter():
        if rPr.tag in (qn('a:rPr'), qn('a:defRPr'), qn('a:endParaRPr')):
            for tag in ('a:latin', 'a:ea', 'a:cs'):
                for old in rPr.findall(qn(tag)):
                    rPr.remove(old)
                rPr.append(rPr.makeelement(qn(tag), {'typeface': FONT}))


for sl in prs.slides:
    for sh in sl.shapes:
        stamp(sh._element)
        if sh.has_chart:
            stamp(sh.chart._chartSpace)

out = 'ppt資料/2分プレゼン_v3.pptx'
prs.save(out)
print('saved:', out, '/ 5 slides')
